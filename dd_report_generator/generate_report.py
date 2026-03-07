"""ExcelテンプレートからPowerPoint DD報告書を自動生成するスクリプト。

使い方:
    python -m dd_report_generator.generate_report -i data.xlsx -o report.pptx
"""

import argparse
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- デザイン定数 ----------
SLIDE_WIDTH = Inches(13.333)  # 16:9 ワイド
SLIDE_HEIGHT = Inches(7.5)

COLOR_PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)    # 濃紺
COLOR_ACCENT = RGBColor(0x44, 0x72, 0xC4)      # 青
COLOR_LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)    # 薄いグレー
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x33, 0x33, 0x33)
COLOR_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_RED = RGBColor(0xC0, 0x39, 0x2B)
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)
COLOR_YELLOW = RGBColor(0xF3, 0x9C, 0x12)
COLOR_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
COLOR_ROW_ALT = RGBColor(0xE8, 0xEE, 0xF5)

FONT_TITLE = "Meiryo"
FONT_BODY = "Meiryo"

CONTENT_LEFT = Inches(0.6)
CONTENT_TOP = Inches(1.4)
CONTENT_WIDTH = Inches(12.1)
CONTENT_HEIGHT = Inches(5.6)


def _read_sheet_data(wb, sheet_name):
    """シートの全データを2次元リストで返す。シートが存在しなければ空リスト。"""
    if sheet_name not in wb.sheetnames:
        return []
    ws = wb[sheet_name]
    data = []
    for row in ws.iter_rows(values_only=True):
        data.append(list(row))
    return data


def _add_slide_title_bar(slide, title_text):
    """スライド上部にタイトルバーを追加。"""
    # 背景バー
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # タイトルテキスト
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_TITLE
    p.font.bold = True


def _add_footer(slide, company_name=""):
    """スライド下部にフッターを追加。"""
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.3),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(7.2), Inches(5), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"CONFIDENTIAL | {company_name}" if company_name else "CONFIDENTIAL"
    p.font.size = Pt(8)
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_BODY


def _add_table(slide, data, left, top, width, height, header_rows=1):
    """データの2次元リストからテーブルを作成。"""
    if not data or not data[0]:
        return None

    rows = len(data)
    cols = len(data[0])
    col_width = int(width / cols)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 列幅を均等割（最後の列で余りを吸収）
    for i in range(cols):
        table.columns[i].width = col_width

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = _format_value(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.name = FONT_BODY

                if r_idx < header_rows:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLOR_WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = COLOR_DARK_TEXT
                    # 数値は右寄せ
                    if isinstance(val, (int, float)):
                        paragraph.alignment = PP_ALIGN.RIGHT
                    else:
                        paragraph.alignment = PP_ALIGN.LEFT

            # セル背景
            if r_idx < header_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_HEADER_BG
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE

    return table_shape


def _format_value(val):
    """セル値を表示用に変換。"""
    if val is None:
        return ""
    if isinstance(val, float):
        if val == int(val):
            return f"{int(val):,}"
        return f"{val:,.1f}"
    if isinstance(val, int):
        return f"{val:,}"
    return str(val)


def _build_cover_slide(prs, data):
    """表紙スライド。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # 全面背景
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    # アクセントライン
    line = slide.shapes.add_shape(
        1, Inches(2), Inches(3.5), Inches(9.333), Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # タイトル
    info = {}
    for row in data[1:]:  # skip header
        if row[0] and row[1]:
            info[row[0]] = row[1]

    title_text = info.get("案件名", "デューデリジェンス報告書")
    target = info.get("対象会社名", "")
    date_str = _format_value(info.get("報告日", ""))
    author = info.get("作成者", "")
    author_company = info.get("作成会社", "")

    txBox = slide.shapes.add_textbox(
        Inches(2), Inches(1.5), Inches(9.333), Inches(1.5),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_TITLE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 対象会社
    if target:
        txBox2 = slide.shapes.add_textbox(
            Inches(2), Inches(2.6), Inches(9.333), Inches(0.8),
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"対象会社: {target}"
        p2.font.size = Pt(22)
        p2.font.color.rgb = COLOR_WHITE
        p2.font.name = FONT_TITLE
        p2.alignment = PP_ALIGN.CENTER

    # 作成情報
    bottom_text = "  |  ".join(filter(None, [author_company, author, date_str]))
    if bottom_text:
        txBox3 = slide.shapes.add_textbox(
            Inches(2), Inches(4.0), Inches(9.333), Inches(0.6),
        )
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = bottom_text
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        p3.font.name = FONT_BODY
        p3.alignment = PP_ALIGN.CENTER


def _build_kv_slide(prs, title, data, company_name):
    """キー・バリュー形式のスライド（会社概要、総括）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title_bar(slide, title)
    _add_footer(slide, company_name)

    # テーブルとして表示（ヘッダー行を含む）
    filtered = [row[:2] for row in data if row[0] is not None]
    if filtered:
        _add_table(
            slide, filtered,
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
            header_rows=1,
        )


def _build_financial_slide(prs, title, data, company_name):
    """BS/PL等の財務テーブルスライド。空行でのセクション分けを維持。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title_bar(slide, title)
    _add_footer(slide, company_name)

    # ヘッダー行 + 期情報行 + 科目行
    # 空行・None行を除外しつつ、セクションヘッダーは維持
    table_data = []
    for row in data:
        if all(v is None or v == "" for v in row):
            continue
        table_data.append(row)

    if not table_data:
        return

    # データが多い場合は複数スライドに分割
    max_rows_per_slide = 25
    header = table_data[0]

    if len(table_data) <= max_rows_per_slide:
        _add_table(
            slide, table_data,
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
            header_rows=1 if len(table_data) > 1 else 0,
        )
    else:
        # 最初のスライドに最初の部分
        chunk = [header] + table_data[1:max_rows_per_slide]
        _add_table(
            slide, chunk,
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
            header_rows=1,
        )
        # 残りを追加スライドへ
        remaining = table_data[max_rows_per_slide:]
        while remaining:
            chunk = [header] + remaining[:max_rows_per_slide - 1]
            remaining = remaining[max_rows_per_slide - 1:]
            extra_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_slide_title_bar(extra_slide, f"{title}（続き）")
            _add_footer(extra_slide, company_name)
            _add_table(
                extra_slide, chunk,
                CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
                header_rows=1,
            )


def _build_verification_slide(prs, data, company_name):
    """科目別検証結果スライド。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title_bar(slide, "科目別検証結果")
    _add_footer(slide, company_name)

    # 入力済み行のみ抽出
    table_data = [data[0]]  # header
    for row in data[1:]:
        if row[1] is not None and str(row[1]).strip():
            table_data.append(row)

    if len(table_data) <= 1:
        # データなし
        txBox = slide.shapes.add_textbox(
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(1),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "（データが入力されていません）"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BLACK
        p.font.name = FONT_BODY
        return

    max_rows_per_slide = 15
    header = table_data[0]

    if len(table_data) <= max_rows_per_slide:
        _add_table(
            slide, table_data,
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
            header_rows=1,
        )
    else:
        chunk = [header] + table_data[1:max_rows_per_slide]
        _add_table(
            slide, chunk,
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
            header_rows=1,
        )
        remaining = table_data[max_rows_per_slide:]
        while remaining:
            chunk = [header] + remaining[:max_rows_per_slide - 1]
            remaining = remaining[max_rows_per_slide - 1:]
            extra_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_slide_title_bar(extra_slide, "科目別検証結果（続き）")
            _add_footer(extra_slide, company_name)
            _add_table(
                extra_slide, chunk,
                CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
                header_rows=1,
            )


def _build_tax_dd_slide(prs, data, company_name):
    """税務DDスライド。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title_bar(slide, "税務デューデリジェンス")
    _add_footer(slide, company_name)

    table_data = [data[0]]
    for row in data[1:]:
        if row[0] is not None and str(row[0]).strip():
            table_data.append(row)

    if len(table_data) > 1:
        _add_table(
            slide, table_data,
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
            header_rows=1,
        )


def _get_company_name(wb):
    """表紙シートから対象会社名を取得。"""
    data = _read_sheet_data(wb, "表紙")
    for row in data[1:]:
        if row[0] == "対象会社名" and row[1]:
            return str(row[1])
    return ""


def generate_report(input_path: str, output_path: str):
    wb = load_workbook(input_path, data_only=True)
    company_name = _get_company_name(wb)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. 表紙
    cover_data = _read_sheet_data(wb, "表紙")
    if cover_data:
        _build_cover_slide(prs, cover_data)

    # 2. 会社概要
    overview_data = _read_sheet_data(wb, "会社概要")
    if overview_data:
        _build_kv_slide(prs, "会社概要", overview_data, company_name)

    # 3. 貸借対照表
    bs_data = _read_sheet_data(wb, "貸借対照表")
    if bs_data:
        _build_financial_slide(prs, "貸借対照表（BS）推移 — 直近3期", bs_data, company_name)

    # 4. 損益計算書
    pl_data = _read_sheet_data(wb, "損益計算書")
    if pl_data:
        _build_financial_slide(prs, "損益計算書（PL）推移 — 直近3期", pl_data, company_name)

    # 5. 修正貸借対照表
    adj_bs_data = _read_sheet_data(wb, "修正貸借対照表")
    if adj_bs_data:
        _build_financial_slide(prs, "修正貸借対照表（時価純資産）", adj_bs_data, company_name)

    # 6. 科目別検証
    verif_data = _read_sheet_data(wb, "科目別検証")
    if verif_data:
        _build_verification_slide(prs, verif_data, company_name)

    # 7. 税務DD
    tax_data = _read_sheet_data(wb, "税務DD")
    if tax_data:
        _build_tax_dd_slide(prs, tax_data, company_name)

    # 8. 総括
    summary_data = _read_sheet_data(wb, "総括")
    if summary_data:
        _build_kv_slide(prs, "総括・リスクサマリー", summary_data, company_name)

    prs.save(output_path)
    print(f"レポートを生成しました: {output_path}")
    print(f"スライド数: {len(prs.slides)}")


def main():
    parser = argparse.ArgumentParser(description="Excel入力データからDD報告書PowerPointを生成")
    parser.add_argument(
        "-i", "--input",
        required=True,
        help="入力Excelファイルパス",
    )
    parser.add_argument(
        "-o", "--output",
        default="output/dd_report.pptx",
        help="出力PowerPointファイルパス (default: output/dd_report.pptx)",
    )
    args = parser.parse_args()
    Path(args.output).parent.mkdir(parents=True, exist_ok=True)
    generate_report(args.input, args.output)


if __name__ == "__main__":
    main()
